"""Generate PowerPoint presentation for the smiles-js-eval study."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ── Colour palette ──────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x11, 0x17)   # near-black
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0x8B, 0x8B, 0x8B)
ACCENT    = RGBColor(0x4E, 0xC9, 0xB0)   # teal/mint
ACCENT2   = RGBColor(0xFF, 0xA0, 0x50)   # orange
RED       = RGBColor(0xFF, 0x55, 0x55)
GREEN     = RGBColor(0x55, 0xFF, 0x88)
BLUE      = RGBColor(0x55, 0x99, 0xFF)
YELLOW    = RGBColor(0xFF, 0xDD, 0x57)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helpers ─────────────────────────────────────────────────────────
def solid_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name="Consolas"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_para(tf, text, size=18, bold=False, color=WHITE, font_name="Consolas", space_before=0, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.space_before = Pt(space_before)
    p.alignment = align
    return p

def add_table(slide, rows, cols, data, left, top, width, height, header_color=ACCENT, cell_color=WHITE, font_size=11):
    """data = list of lists; first row is header."""
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Consolas"
                paragraph.font.color.rgb = header_color if i == 0 else cell_color
                paragraph.font.bold = (i == 0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1A, 0x1E, 0x24) if i == 0 else RGBColor(0x12, 0x16, 0x1C)
    return table

# ── SLIDE 1: Title ──────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
solid_bg(slide)
add_text(slide, 1, 1.5, 11, 1.2, "High Heels for Haiku", size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, 1, 2.8, 11, 0.8, "When does code representation help LLMs", size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.4, 11, 0.8, "reason about molecules?", size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 1, 4.8, 11, 0.5, "smiles-js-eval  ·  15 qualitative probes  ·  Haiku", size=16, color=GREY, align=PP_ALIGN.CENTER)
add_text(slide, 1, 5.5, 11, 0.5, "February 2026", size=14, color=GREY, align=PP_ALIGN.CENTER)

# ── SLIDE 2: The Question ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.7, "The Question", size=32, bold=True, color=ACCENT)
tf = add_text(slide, 0.8, 1.3, 11.5, 1, "Can we help LLMs reason about chemistry by converting", size=22, color=WHITE)
add_para(tf, "SMILES strings into composable JavaScript code?", size=22, color=YELLOW, bold=True, space_before=4)
add_para(tf, "", size=12)
add_para(tf, "SMILES (string):", size=16, color=GREY, space_before=16)
add_para(tf, "  CNc1ncnc2c1ncn2Cc1ccccc1", size=18, color=RED)
add_para(tf, "", size=8)
add_para(tf, "smiles-js Code (tree):", size=16, color=GREY, space_before=8)
add_para(tf, "  pyrimidine = Fragment('c1ncncc1')", size=18, color=GREEN)
add_para(tf, "  imidazole  = Fragment('c2cncn2')", size=18, color=GREEN)
add_para(tf, "  purine     = pyrimidine.fuse(4, imidazole)", size=18, color=GREEN)
add_para(tf, "  benzene    = Fragment('c1ccccc1')", size=18, color=GREEN)
add_para(tf, "  molecule   = Molecule([purine, benzene])", size=18, color=GREEN)

# ── SLIDE 3: Study Design ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.7, "Study Design", size=32, bold=True, color=ACCENT)
tf = add_text(slide, 0.8, 1.3, 5.5, 5, "Method", size=20, bold=True, color=ACCENT2)
add_para(tf, "", size=6)
add_para(tf, "• Same molecule, same question", size=17, color=WHITE, space_before=8)
add_para(tf, "• Give Haiku the SMILES string", size=17, color=WHITE, space_before=4)
add_para(tf, "• Give Haiku the smiles-js code", size=17, color=WHITE, space_before=4)
add_para(tf, "• Compare answers against ground truth", size=17, color=WHITE, space_before=4)
add_para(tf, "• 15 probes × 6 task types", size=17, color=WHITE, space_before=4)

tf2 = add_text(slide, 7, 1.3, 5.5, 5, "Task Types", size=20, bold=True, color=ACCENT2)
add_para(tf2, "", size=6)
add_para(tf2, "1. Aromatic ring counting", size=17, color=WHITE, space_before=8)
add_para(tf2, "2. Scaffold recognition", size=17, color=WHITE, space_before=4)
add_para(tf2, "3. Functional group detection", size=17, color=WHITE, space_before=4)
add_para(tf2, "4. H-bond donor/acceptor counting", size=17, color=WHITE, space_before=4)
add_para(tf2, "5. Stereocenter counting", size=17, color=WHITE, space_before=4)
add_para(tf2, "6. Ring size discrimination", size=17, color=WHITE, space_before=4)

# ── SLIDE 4: Scorecard ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.7, "Scorecard: Code Wins 7 – SMILES Wins 4 – Tie 4", size=30, bold=True, color=ACCENT)

data = [
    ["#", "Task", "Molecule", "Truth", "SMILES", "Code", "Winner"],
    ["1", "Ring count", "Purine + 2 benzenes", "4", "3 ✗", "4 ✓", "Code"],
    ["2", "Ring count", "Benzodifuran", "3", "2 ✗", "3 ✓", "Code"],
    ["3", "Ring count", "Coumarin-furan", "3", "2 ✗", "3 ✓", "Code"],
    ["4", "Ring count", "Purine + benzene", "3", "3 ✓", "4 ✗", "SMILES"],
    ["5", "Ring + relabel", "Same as #4", "3", "—", "3 ✓", "Code"],
    ["6", "Scaffold", "Diaminoquinazoline", "quinazoline", "benzimidazole ✗", "quinazoline ✓", "Code"],
    ["7", "Func groups", "4-group molecule", "4 groups", "4/4 ✓", "3/4 + 2 FP", "SMILES"],
    ["8", "Mol ID", "Ergot alkaloid", "bromocriptine", "indinavir ✗", "indinavir ✗", "Tie"],
    ["9", "Ring count", "1 arom / 2 non-arom", "1", "2 ✗", "1 ✓", "Code"],
    ["10", "Scaffold", "Bis-indolyl uracil", "bis-indolyl", "indazolo ✗", "bis-indolyl ✓", "Code"],
    ["11", "Stereocenters", "Beta-lactam", "3", "3 ✓", "hallucinated", "SMILES"],
    ["12", "H-bond count", "Pentapeptide", "hbd=7 hba=6", "hbd=12 ✗", "7,6 ✓", "Code"],
    ["13", "Scaffold", "Imidazole hub", "imidazole", "✓ (rambling)", "✓ (crisp)", "Tie"],
    ["14", "Amide count", "Pentapeptide", "5", "5 ✓", "3 ✗", "SMILES"],
    ["15", "Ring size", "Mixed 5/6-mem", "4 rings", "✓", "✓", "Tie"],
]
add_table(slide, 16, 7, data, 0.3, 1.1, 12.7, 6.2, font_size=10)

# ── SLIDE 5: Pattern 1 – Fused Ring Counting ────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.7, "Pattern 1: Fused Ring Counting", size=32, bold=True, color=ACCENT)
add_text(slide, 0.8, 1.0, 11, 0.5, "Code wins 4/5 hard cases  vs  SMILES wins 2/5", size=18, color=ACCENT2)

tf = add_text(slide, 0.8, 1.8, 5.5, 4.5, "The Problem with SMILES", size=20, bold=True, color=RED)
add_para(tf, "", size=6)
add_para(tf, "SMILES compresses fused rings:", size=16, color=WHITE, space_before=8)
add_para(tf, "  c1c2ccoc2cc2ccc(=O)oc12", size=16, color=RED, space_before=4)
add_para(tf, "", size=4)
add_para(tf, "Ring closure digits (c1...c12) are", size=16, color=WHITE, space_before=8)
add_para(tf, "nested and hard to parse.", size=16, color=WHITE)
add_para(tf, "Haiku sees 2 rings. Truth: 3.", size=16, color=YELLOW, space_before=8)

tf2 = add_text(slide, 7, 1.8, 5.5, 4.5, "How Code Helps", size=20, bold=True, color=GREEN)
add_para(tf2, "", size=6)
add_para(tf2, "Code enumerates each ring:", size=16, color=WHITE, space_before=8)
add_para(tf2, "  ring1 = Fragment('c1cccoc1')", size=14, color=GREEN, space_before=4)
add_para(tf2, "  ring2 = Fragment('c2occc2')", size=14, color=GREEN, space_before=2)
add_para(tf2, "  ring3 = Fragment('c2cccoc2')", size=14, color=GREEN, space_before=2)
add_para(tf2, "  FusedRing([ring1, ring2, ring3])", size=14, color=GREEN, space_before=2)
add_para(tf2, "", size=4)
add_para(tf2, "3 Fragments in FusedRing = 3 rings.", size=16, color=YELLOW, space_before=8)
add_para(tf2, "Counting nodes in a tree is easy.", size=16, color=WHITE, space_before=4)

# ── SLIDE 6: Pattern 2 – Scaffold Recognition ───────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.7, "Pattern 2: Scaffold Recognition", size=32, bold=True, color=ACCENT)
add_text(slide, 0.8, 1.0, 11, 0.5, "Code correctly names scaffolds that SMILES obscures", size=18, color=ACCENT2)

tf = add_text(slide, 0.8, 1.8, 5.5, 4.5, "SMILES → Wrong Answer", size=20, bold=True, color=RED)
add_para(tf, "", size=6)
add_para(tf, "Nc1nc(N)c2c(...)cccc2n1", size=15, color=RED, space_before=8)
add_para(tf, "", size=4)
add_para(tf, 'Haiku says: "benzimidazole"', size=16, color=WHITE, space_before=8)
add_para(tf, "Truth: quinazoline", size=16, color=YELLOW, space_before=4)
add_para(tf, "", size=8)
add_para(tf, "The fused ring looks like one blob.", size=16, color=GREY, space_before=8)
add_para(tf, "Haiku can't tell pyrimidine+benzene", size=16, color=GREY, space_before=2)
add_para(tf, "from imidazole+benzene.", size=16, color=GREY, space_before=2)

tf2 = add_text(slide, 7, 1.8, 5.5, 4.5, "Code → Correct Answer", size=20, bold=True, color=GREEN)
add_para(tf2, "", size=6)
add_para(tf2, "diaminopyrimidine = Fragment(...)", size=15, color=GREEN, space_before=8)
add_para(tf2, "benzene = Fragment('c1ccccc1')", size=15, color=GREEN, space_before=2)
add_para(tf2, "core = diaminopyrimidine.fuse(benzene)", size=15, color=GREEN, space_before=2)
add_para(tf2, "", size=4)
add_para(tf2, 'Haiku says: "quinazoline" ✓', size=16, color=WHITE, space_before=8)
add_para(tf2, "", size=8)
add_para(tf2, "pyrimidine.fuse(benzene) is the", size=16, color=GREY, space_before=8)
add_para(tf2, "textbook definition of quinazoline.", size=16, color=GREY, space_before=2)
add_para(tf2, "Code makes the recipe visible.", size=16, color=GREY, space_before=2)

# ── SLIDE 7: Pattern 3 – Repeating Units ────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.7, "Pattern 3: Counting in Long Chains", size=32, bold=True, color=ACCENT)
add_text(slide, 0.8, 1.0, 11, 0.5, "Peptide H-bond counting: Code nails it, SMILES overcounts", size=18, color=ACCENT2)

tf = add_text(slide, 0.8, 1.8, 5.5, 4.5, "SMILES: Wall of Text", size=20, bold=True, color=RED)
add_para(tf, "", size=6)
add_para(tf, "Cc1cccc(NC(=O)N[C@@H](Cc2", size=14, color=RED, space_before=8)
add_para(tf, "ccccc2)C(=O)N[C@H](CC(C)C", size=14, color=RED, space_before=2)
add_para(tf, ")C(=O)N[C@@H](...) ...", size=14, color=RED, space_before=2)
add_para(tf, "", size=4)
add_para(tf, "Haiku says: hbd = 12", size=16, color=WHITE, space_before=8)
add_para(tf, "Truth:      hbd = 7", size=16, color=YELLOW, space_before=4)
add_para(tf, "", size=4)
add_para(tf, "Double-counted NH donors in the", size=16, color=GREY, space_before=8)
add_para(tf, "repetitive peptide backbone.", size=16, color=GREY, space_before=2)

tf2 = add_text(slide, 7, 1.8, 5.5, 4.5, "Code: Structured Decomposition", size=20, bold=True, color=GREEN)
add_para(tf2, "", size=6)
add_para(tf2, "backbone = Fragment('NCN...NCN...')", size=14, color=GREEN, space_before=8)
add_para(tf2, "carbonyl1 = Linear(['O'], ['='])", size=14, color=GREEN, space_before=2)
add_para(tf2, "carbonyl2 = Linear(['O'], ['='])", size=14, color=GREEN, space_before=2)
add_para(tf2, "...6 carbonyls total...", size=14, color=GREEN, space_before=2)
add_para(tf2, "", size=4)
add_para(tf2, "Haiku says: hbd = 7, hba = 6  ✓", size=16, color=WHITE, space_before=8)
add_para(tf2, "", size=4)
add_para(tf2, "Each carbonyl is a discrete countable", size=16, color=GREY, space_before=8)
add_para(tf2, "node. No double-counting possible.", size=16, color=GREY, space_before=2)

# ── SLIDE 8: When SMILES Wins ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.7, "When SMILES Wins (4 cases)", size=32, bold=True, color=ACCENT)

tf = add_text(slide, 0.8, 1.3, 5.5, 5, "Pattern A: Functional Groups", size=20, bold=True, color=RED)
add_para(tf, "", size=6)
add_para(tf, "SMILES keeps groups inline:", size=16, color=WHITE, space_before=8)
add_para(tf, "  C(=O)N  → amide  (obvious)", size=15, color=YELLOW, space_before=4)
add_para(tf, "  C(=O)O  → carboxyl", size=15, color=YELLOW, space_before=2)
add_para(tf, "  [N+](=O)[O-] → nitro", size=15, color=YELLOW, space_before=2)
add_para(tf, "", size=4)
add_para(tf, "Code splits these across fragments:", size=16, color=WHITE, space_before=8)
add_para(tf, "  Fragment('NC') → looks like methylamine", size=15, color=RED, space_before=4)
add_para(tf, "  Linear(['O'],['=']) → could be anything", size=15, color=RED, space_before=2)
add_para(tf, "", size=4)
add_para(tf, "→ Code hallucinated amine, ester", size=16, color=GREY, space_before=8)

tf2 = add_text(slide, 7, 1.3, 5.5, 5, "Pattern B: Atom-Level Annotations", size=20, bold=True, color=RED)
add_para(tf2, "", size=6)
add_para(tf2, "Stereocenters: [C@H], [C@@H]", size=16, color=WHITE, space_before=8)
add_para(tf2, "These are SMILES annotations.", size=16, color=WHITE, space_before=4)
add_para(tf2, "", size=4)
add_para(tf2, "Code fragments still contain the", size=16, color=GREY, space_before=8)
add_para(tf2, "same [C@H] notation internally.", size=16, color=GREY, space_before=2)
add_para(tf2, "No structural advantage.", size=16, color=GREY, space_before=2)
add_para(tf2, "", size=4)
add_para(tf2, "SMILES: counted 3 stereocenters ✓", size=16, color=GREEN, space_before=8)
add_para(tf2, "Code:   hallucinated 1425 ✗", size=16, color=RED, space_before=4)

# ── SLIDE 9: The Fundamental Tradeoff ────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.7, "The Fundamental Tradeoff", size=32, bold=True, color=ACCENT)

tf = add_text(slide, 1.5, 1.5, 10, 1.5, "Code is a tree decomposition of a string.", size=28, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
add_para(tf, "", size=12)
add_para(tf, "Trees are better for counting nodes.", size=24, color=GREEN, space_before=12, align=PP_ALIGN.CENTER)
add_para(tf, "(rings, repeating units, building blocks)", size=18, color=GREY, space_before=4, align=PP_ALIGN.CENTER)
add_para(tf, "", size=12)
add_para(tf, "Strings are better for pattern matching.", size=24, color=RED, space_before=12, align=PP_ALIGN.CENTER)
add_para(tf, "(functional groups, stereocenters, bond patterns)", size=18, color=GREY, space_before=4, align=PP_ALIGN.CENTER)

# ── SLIDE 10: Sonnet 4.5 Batch Results (context) ────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.7, "Earlier: Sonnet 4.5 Batch Results (n=100)", size=30, bold=True, color=ACCENT)
add_text(slide, 0.8, 1.0, 11, 0.5, "Quantitative baseline before the qualitative probing", size=18, color=GREY)

data = [
    ["Task", "SMILES", "Code", "Code+Relabel", "Winner"],
    ["Ring Count", "92%", "76%", "74%", "SMILES"],
    ["BBBP (bal. acc)", "59%", "54%", "62%", "Code+Relabel"],
    ["SMILES Repair", "28%", "0%", "0%", "SMILES"],
]
add_table(slide, 4, 5, data, 1.5, 1.8, 10, 2.2, font_size=16)

tf = add_text(slide, 0.8, 4.3, 11, 2.5, "Key takeaway from batch run:", size=18, bold=True, color=ACCENT2)
add_para(tf, "", size=4)
add_para(tf, "SMILES won 2/3 tasks. Hypothesis not broadly supported.", size=17, color=WHITE, space_before=8)
add_para(tf, "But this hid the nuance: code helps on specific sub-problems.", size=17, color=WHITE, space_before=4)
add_para(tf, "The qualitative probing found where those sub-problems live.", size=17, color=YELLOW, space_before=4)

# ── SLIDE 11: Implications ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.7, "Implications for smiles-js", size=32, bold=True, color=ACCENT)

tf = add_text(slide, 0.8, 1.3, 11, 5.5, "", size=18, color=WHITE)
p = tf.paragraphs[0]
p.text = "1. FusedRing() is the killer feature"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = GREEN
p.font.name = "Consolas"
add_para(tf, "   Primary driver of code advantage — explicitly enumerates rings", size=16, color=WHITE, space_before=4)

add_para(tf, "", size=8)
add_para(tf, "2. Fragment granularity matters", size=22, bold=True, color=YELLOW, space_before=12)
add_para(tf, "   Splitting C(=O)O → Fragment('CO') + Linear(['O'],['='])", size=16, color=WHITE, space_before=4)
add_para(tf, "   destroys the carboxyl pattern. Preserve functional group boundaries.", size=16, color=WHITE, space_before=2)

add_para(tf, "", size=8)
add_para(tf, "3. Code+Relabel is underrated", size=22, bold=True, color=BLUE, space_before=12)
add_para(tf, "   Forcing semantic naming (pyrimidineRing, benzeneRing) before", size=16, color=WHITE, space_before=4)
add_para(tf, "   reasoning fixed code's overcounting problem. Worth testing at scale.", size=16, color=WHITE, space_before=2)

add_para(tf, "", size=8)
add_para(tf, "4. Best of both worlds?", size=22, bold=True, color=ACCENT, space_before=12)
add_para(tf, "   Hybrid: SMILES for atom-level patterns + code tree for topology.", size=16, color=WHITE, space_before=4)
add_para(tf, "   Or: larger fragments that preserve functional groups.", size=16, color=WHITE, space_before=2)

# ── SLIDE 12: Closing ───────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide)
add_text(slide, 1, 2.0, 11, 1.2, '"High Heels"', size=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.3, 11, 1, "Code doesn't make Haiku taller.", size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.9, 11, 1, "It gives it a boost exactly where it needs one:", size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 1, 4.7, 11, 0.8, "counting structural units in compressed notation.", size=24, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

tf = add_text(slide, 2, 5.8, 9, 1, "github.com/Ghost---Shadow/smiles-js-eval", size=14, color=GREY, align=PP_ALIGN.CENTER)
add_para(tf, "data/haiku-probes.json  ·  docs/observations.md", size=12, color=GREY, space_before=4, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────
out = "docs/high-heels-study.pptx"
prs.save(out)
print(f"Saved: {out}")
